"""Regression tests for output-language safety."""

import json
import re
import subprocess
import sys
import tempfile
import unittest
from pathlib import Path

from docx import Document
from pptx import Presentation
from pptx.util import Inches

ROOT = Path(__file__).resolve().parents[2]
SCRIPTS = ROOT / "scripts"
sys.path.insert(0, str(SCRIPTS))

import translate_pptx_native

HANGUL = re.compile(r"[가-힣ㄱ-ㅎㅏ-ㅣ]")


def document_text(path):
    """Collect visible paragraph and table-cell text from a DOCX."""
    document = Document(path)
    parts = [paragraph.text for paragraph in document.paragraphs]
    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                parts.extend(paragraph.text for paragraph in cell.paragraphs)
    return "\n".join(parts)


class DocxLanguageSafetyTests(unittest.TestCase):
    def convert(self, directory, language=None):
        """Convert an English fixture whose filename has the former Korean suffix."""
        markdown = directory / "english-report-ko.md"
        output = directory / ("explicit-ko.docx" if language == "ko" else "default.docx")
        markdown.write_text(
            "# English report\n\n"
            "> **Key Takeaways**\n"
            "> - This document must remain in English.\n"
        )
        command = [
            sys.executable,
            str(SCRIPTS / "generate_styled_docx.py"),
            str(markdown),
            "-o",
            str(output),
        ]
        if language:
            command.extend(["-l", language])
        result = subprocess.run(command, capture_output=True, text=True, check=False)
        self.assertEqual(result.returncode, 0, result.stderr or result.stdout)
        self.assertTrue(output.exists())
        return document_text(output)

    def test_ko_filename_does_not_enable_korean_mode(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            text = self.convert(Path(temp_dir))
        self.assertIn("Key Takeaways", text)
        self.assertIsNone(HANGUL.search(text), text)

    def test_korean_labels_require_explicit_language(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            text = self.convert(Path(temp_dir), language="ko")
        self.assertIn("핵심 요약", text)


class PptxLanguageSafetyTests(unittest.TestCase):
    def create_mixed_script_deck(self, path):
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        box.text_frame.text = "Seoul (서울)"
        presentation.save(path)

    def test_mixed_script_name_is_an_issue_by_default(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            deck = Path(temp_dir) / "mixed-name.pptx"
            self.create_mixed_script_deck(deck)
            result = translate_pptx_native.review(deck, source_lang="ko")
        self.assertEqual(len(result["issues"]), 1)
        self.assertEqual(result["intentional"], [])

    def test_mixed_script_name_exemption_requires_opt_in(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            deck = Path(temp_dir) / "mixed-name.pptx"
            self.create_mixed_script_deck(deck)
            result = translate_pptx_native.review(
                deck,
                source_lang="ko",
                allow_source_script=True,
            )
        self.assertEqual(result["issues"], [])
        self.assertEqual(len(result["intentional"]), 1)


    def test_native_cli_extract_apply_and_review(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            directory = Path(temp_dir)
            source = directory / "source.pptx"
            extraction = directory / "extract.json"
            translations = directory / "translations.json"
            output = directory / "output.pptx"
            review = directory / "review.json"
            self.create_mixed_script_deck(source)

            helper = SCRIPTS / "translate_pptx_native.py"
            extract_result = subprocess.run(
                [
                    sys.executable,
                    str(helper),
                    "extract",
                    str(source),
                    "--source-lang",
                    "ko",
                    "--output",
                    str(extraction),
                ],
                capture_output=True,
                text=True,
                check=False,
            )
            self.assertEqual(extract_result.returncode, 0, extract_result.stderr)
            data = json.loads(extraction.read_text())
            self.assertEqual(data["stats"]["total_paragraphs"], 1)
            data["slides"][0]["paragraphs"][0]["translated"] = "Seoul"
            translations.write_text(json.dumps(data))

            apply_result = subprocess.run(
                [
                    sys.executable,
                    str(helper),
                    "apply",
                    str(source),
                    "--translations",
                    str(translations),
                    "--output",
                    str(output),
                    "--target-lang",
                    "en",
                ],
                capture_output=True,
                text=True,
                check=False,
            )
            self.assertEqual(apply_result.returncode, 0, apply_result.stderr)
            self.assertTrue(output.exists())

            review_result = subprocess.run(
                [
                    sys.executable,
                    str(helper),
                    "review",
                    str(output),
                    "--source-lang",
                    "ko",
                    "--output",
                    str(review),
                ],
                capture_output=True,
                text=True,
                check=False,
            )
            self.assertEqual(review_result.returncode, 0, review_result.stderr)
            self.assertEqual(json.loads(review.read_text())["issues"], [])

class InstructionSafetyTests(unittest.TestCase):
    def test_every_skill_variant_has_output_language_guardrail(self):
        skill_files = sorted((ROOT / "skills").glob("*/SKILL.md"))
        skill_files += sorted((ROOT / "quick").glob("*/SKILL.md"))
        self.assertEqual(len(skill_files), 8)
        for path in skill_files:
            with self.subTest(path=path):
                text = path.read_text()
                self.assertIsNone(
                    HANGUL.search(text),
                    f"Runtime skill instructions must be English-only: {path}",
                )
                self.assertIn("## Output Language Safety", text)
                self.assertTrue(
                    "Never introduce Korean into non-Korean output" in text
                    or "When the target is not Korean, never generate or retain Korean text" in text
                )

    def test_cli_retention_is_opt_in_and_summary_uses_target_language(self):
        source = (SCRIPTS / "translate_pptx_cli.py").read_text()
        self.assertIn("PRESERVE_SOURCE_SCRIPT = False", source)
        self.assertIn("--preserve-source-script", source)
        self.assertIn("do not retain source-script characters", source)
        self.assertIn("summary entirely in {tgt_name}", source)
        self.assertIn("Do not introduce Korean unless {tgt_name} is Korean", source)
        self.assertNotIn("return texts  # fallback: return originals", source)
        self.assertIn("refusing to reuse source text", source)
        self.assertIn("Refusing to save non-Korean output", source)

    def test_release_install_context_filters_korean(self):
        source = (SCRIPTS / "build-release.sh").read_text()
        self.assertIn("copy_english_install_guide", source)
        self.assertIn("if not hangul.search(line)", source)
        self.assertIn('copy_english_install_guide "$skill_staging/INSTALL.md"', source)
        self.assertIn('copy_english_install_guide "$ALL_STAGING/INSTALL.md"', source)


if __name__ == "__main__":
    unittest.main()
