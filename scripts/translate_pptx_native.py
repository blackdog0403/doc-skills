"""
translate_pptx_native.py — PPTX translation helper for Amazon Quick

This script handles the non-LLM parts of PPTX translation:
  1. extract_texts(pptx_path) → JSON of all translatable paragraphs
  2. apply_translations(pptx_path, translations_json, output_path) → applies translations + font normalization
  3. review(pptx_path, source_lang) → scans for remaining source language text

The agent handles the actual translation (it IS the LLM).
"""

import re
import os
from pptx import Presentation

# ── Font Config ──
FONT_MAP = {
    'ko': '맑은 고딕',
    'ja': 'Yu Gothic UI',
    'en': 'Amazon Ember',
    'zh': 'Microsoft YaHei',
}
FONT_DEFAULT = 'Arial'

FONT_REPLACE = {
    'NanumSquareOTF': '맑은 고딕',
    'NanumSquareOTF Bold': '맑은 고딕',
    'NanumSquareOTF ExtraBold': '맑은 고딕',
    'NanumSquareOTF Light': '맑은 고딕',
}

LANG_PATTERNS = {
    'ko': '[\uac00-\ud7a3]',
    'ja': '[\u3041-\u3093\u30a1-\u30f6\u30fc]',
    'zh': '[\u4e00-\u9fff]',
    'en': '[A-Za-z]',
}

LANG_NAMES = {
    'ko': 'Korean', 'ja': 'Japanese', 'zh': 'Chinese',
    'en': 'English', 'es': 'Spanish', 'fr': 'French',
    'de': 'German', 'pt': 'Portuguese', 'it': 'Italian',
}


def _has_lang(text, lang):
    pattern = LANG_PATTERNS.get(lang)
    if pattern:
        return bool(re.search(pattern, text))
    return bool(text.strip())


def _detect_language(text):
    if not text or not text.strip():
        return None
    counts = {
        'ko': len(re.findall('[\uac00-\ud7a3]', text)),
        'ja': len(re.findall('[\u3041-\u3093\u30a1-\u30f6]', text)),
        'zh': len(re.findall('[\u4e00-\u9fff]', text)),
        'en': len(re.findall('[A-Za-z]', text)),
    }
    dominant = max(counts, key=counts.get)
    return dominant if counts[dominant] > 0 else None


def _iter_all_frames(slide):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    def _from_shape(shape):
        if shape.has_text_frame:
            yield shape.text_frame
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    yield cell.text_frame
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                yield from _from_shape(child)
    for shape in slide.shapes:
        yield from _from_shape(shape)


# ── Step 1: Extract ──

def extract_texts(pptx_path, source_lang='ko', include_notes=True):
    """
    Extract all translatable paragraphs from a PPTX.
    Returns a structure:
    {
        "source_lang": "ko",
        "total_slides": 30,
        "slides": [
            {
                "slide_num": 1,
                "paragraphs": [
                    {"frame_idx": 0, "para_idx": 0, "text": "한국어 텍스트"},
                    ...
                ]
            },
            ...
        ],
        "stats": {"total_paragraphs": 150, "total_chars": 12000, "slides_with_content": 25}
    }
    """
    prs = Presentation(pptx_path)
    result = {
        "source_lang": source_lang,
        "total_slides": len(prs.slides),
        "slides": [],
        "stats": {"total_paragraphs": 0, "total_chars": 0, "slides_with_content": 0}
    }

    for si, slide in enumerate(prs.slides):
        slide_data = {"slide_num": si + 1, "paragraphs": []}
        frames = list(_iter_all_frames(slide))

        # Add notes frame
        if include_notes and slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            frames.append(slide.notes_slide.notes_text_frame)

        for fi, frame in enumerate(frames):
            for pi, para in enumerate(frame.paragraphs):
                text = para.text.strip()
                if text and _has_lang(text, source_lang):
                    slide_data["paragraphs"].append({
                        "frame_idx": fi,
                        "para_idx": pi,
                        "text": para.text  # preserve original whitespace
                    })
                    result["stats"]["total_paragraphs"] += 1
                    result["stats"]["total_chars"] += len(para.text)

        if slide_data["paragraphs"]:
            result["stats"]["slides_with_content"] += 1
        result["slides"].append(slide_data)

    return result


# ── Step 2: Apply Translations ──

def apply_translations(pptx_path, translations, output_path, target_lang='en'):
    """
    Apply translations to a PPTX file.

    translations: list of dicts matching extract_texts output, with added "translated" key:
    [
        {"slide_num": 1, "paragraphs": [
            {"frame_idx": 0, "para_idx": 0, "text": "원본", "translated": "Translated"},
            ...
        ]},
        ...
    ]
    """
    prs = Presentation(pptx_path)

    applied = 0
    for slide_data in translations:
        si = slide_data["slide_num"] - 1
        slide = prs.slides[si]
        frames = list(_iter_all_frames(slide))

        # Add notes frame
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            frames.append(slide.notes_slide.notes_text_frame)

        for para_data in slide_data.get("paragraphs", []):
            fi = para_data["frame_idx"]
            pi = para_data["para_idx"]
            translated = para_data.get("translated", "")
            if not translated:
                continue

            if fi < len(frames):
                frame = frames[fi]
                if pi < len(frame.paragraphs):
                    para = frame.paragraphs[pi]
                    if para.runs:
                        para.runs[0].text = translated
                        for run in para.runs[1:]:
                            run.text = ""
                        applied += 1

    # Post-processing
    artifact_count = _clean_artifacts(prs)
    run_count, xml_count = _normalize_fonts(prs, target_lang)

    # Save with docProps patching
    prs.save(output_path)
    _patch_docprops(output_path)

    return {
        "paragraphs_applied": applied,
        "artifacts_cleaned": artifact_count,
        "fonts_normalized": run_count,
        "xml_fonts_patched": xml_count,
        "output_path": output_path
    }


# ── Step 3: Review ──

def review(pptx_path, source_lang='ko', include_notes=True):
    """Scan for remaining source language text after translation."""
    prs = Presentation(pptx_path)
    remaining = []

    for si, slide in enumerate(prs.slides):
        for frame in _iter_all_frames(slide):
            for para in frame.paragraphs:
                if para.text.strip() and _has_lang(para.text, source_lang):
                    # Check if it's intentional (e.g. "Seoul (서울)")
                    intentional = bool(re.search(
                        r'[A-Za-z].*\([^)]*[\uac00-\ud7a3\u3041-\u30f6\u4e00-\u9fff]|'
                        r'[\uac00-\ud7a3\u3041-\u30f6\u4e00-\u9fff].*\(.*[A-Za-z]',
                        para.text
                    ))
                    remaining.append({
                        "slide": si + 1,
                        "location": "text",
                        "text": para.text[:120],
                        "intentional": intentional
                    })

        if include_notes and slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            for para in slide.notes_slide.notes_text_frame.paragraphs:
                if para.text.strip() and _has_lang(para.text, source_lang):
                    intentional = bool(re.search(
                        r'[A-Za-z].*\([^)]*[\uac00-\ud7a3\u3041-\u30f6\u4e00-\u9fff]|'
                        r'[\uac00-\ud7a3\u3041-\u30f6\u4e00-\u9fff].*\(.*[A-Za-z]',
                        para.text
                    ))
                    remaining.append({
                        "slide": si + 1,
                        "location": "notes",
                        "text": para.text[:120],
                        "intentional": intentional
                    })

    issues = [r for r in remaining if not r["intentional"]]
    intentional = [r for r in remaining if r["intentional"]]
    return {"issues": issues, "intentional": intentional, "total_remaining": len(remaining)}


# ── Internal helpers ──

def _clean_artifacts(prs):
    count = 0
    for slide in prs.slides:
        frames = list(_iter_all_frames(slide))
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            frames.append(slide.notes_slide.notes_text_frame)
        for frame in frames:
            for para in frame.paragraphs:
                for run in para.runs:
                    if '_x000B_' in run.text:
                        run.text = run.text.replace('_x000B_', '\n')
                        count += 1
    return count


def _normalize_fonts(prs, target_lang='en'):
    run_count = 0
    for slide in prs.slides:
        frames = list(_iter_all_frames(slide))
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            frames.append(slide.notes_slide.notes_text_frame)
        for frame in frames:
            for para in frame.paragraphs:
                for run in para.runs:
                    if not run.text.strip():
                        continue
                    lang = _detect_language(run.text)
                    if lang:
                        run.font.name = FONT_MAP.get(lang, FONT_DEFAULT)
                        run_count += 1

    # XML-level font patching
    xml_count = 0
    parts = []
    for slide in prs.slides:
        parts.append(slide._element)
        if slide.has_notes_slide:
            parts.append(slide.notes_slide._element)
    for layout in prs.slide_layouts:
        parts.append(layout._element)
    for master in prs.slide_masters:
        parts.append(master._element)

    for part in parts:
        for elem in part.iter():
            typeface = elem.get('typeface')
            if typeface and typeface in FONT_REPLACE:
                elem.set('typeface', FONT_REPLACE[typeface])
                xml_count += 1

    return run_count, xml_count


def _patch_docprops(pptx_path):
    """Patch docProps/app.xml to replace font references in the ZIP."""
    import zipfile
    import shutil
    tmp = pptx_path + '.tmp'
    try:
        with zipfile.ZipFile(pptx_path, 'r') as zin, zipfile.ZipFile(tmp, 'w') as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename == 'docProps/app.xml':
                    text = data.decode('utf-8')
                    for old_font, new_font in FONT_REPLACE.items():
                        text = text.replace(old_font, new_font)
                    data = text.encode('utf-8')
                zout.writestr(item, data)
        shutil.move(tmp, pptx_path)
    except Exception:
        if os.path.exists(tmp):
            os.remove(tmp)
