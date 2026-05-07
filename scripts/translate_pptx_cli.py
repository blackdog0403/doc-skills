#!/Users/kwangyou/.local/share/translate-pptx-venv/bin/python
"""
PPTX Translator using AWS Bedrock (Claude Opus)

Key design decisions:
- Paragraph-level translation (not run-level) to preserve sentence context
- Translated text goes into first run, remaining runs cleared (preserves first run's formatting)
- Batch API calls (up to 20 paragraphs per request) for efficiency
- Post-translation review pass catches any remaining source language text
- Real-time progress display with slide-by-slide status
- Retry logic for API failures
"""

import json, re, sys, time, argparse
from pptx import Presentation
from pptx.util import Pt
import boto3

# ── Config ──
MODEL_ID = "us.anthropic.claude-opus-4-6-v1"
REGION = "us-west-2"
BATCH_SIZE = 20
MAX_RETRIES = 3

# Font rules by detected language
FONT_KOREAN = "맑은 고딕"
FONT_JAPANESE = "Yu Gothic UI"
FONT_ENGLISH = "Amazon Ember"
FONT_CHINESE = "Microsoft YaHei"
FONT_DEFAULT = "Arial"

# Font replacement map for missing fonts
FONT_REPLACE = {
    'NanumSquareOTF': FONT_KOREAN,
    'NanumSquareOTF Bold': FONT_KOREAN,
    'NanumSquareOTF ExtraBold': FONT_KOREAN,
    'NanumSquareOTF Light': FONT_KOREAN,
}

client = boto3.client('bedrock-runtime', region_name=REGION)

# Language detection patterns
LANG_PATTERNS = {
    'ko': '[가-힣]',
    'ja': '[ぁ-んァ-ヶー]',
    'zh': '[一-鿿]',
    'en': '[A-Za-z]',
}

LANG_NAMES = {
    'ko': 'Korean', 'ja': 'Japanese', 'zh': 'Chinese',
    'en': 'English', 'es': 'Spanish', 'fr': 'French',
    'de': 'German', 'pt': 'Portuguese', 'it': 'Italian',
}

# Runtime state (set by main)
SOURCE_LANG = 'ko'
TARGET_LANG = 'en'
GLOSSARY = {}  # term → translation mapping, loaded from --glossary file
TOTAL_INPUT_TOKENS = 0
TOTAL_OUTPUT_TOKENS = 0


def has_source_lang(text):
    """Check if text contains source language characters."""
    pattern = LANG_PATTERNS.get(SOURCE_LANG)
    if pattern:
        return bool(re.search(pattern, text))
    # For languages without character-based detection, assume all non-target text needs translation
    return bool(text.strip())


def has_korean(text):
    return bool(re.search('[가-힣]', text))


def detect_language(text):
    """Detect dominant language in text for font selection."""
    if not text or not text.strip():
        return None
    korean = len(re.findall('[가-힣]', text))
    japanese = len(re.findall('[ぁ-んァ-ヶ]', text))
    chinese = len(re.findall('[一-鿿]', text))
    latin = len(re.findall('[A-Za-z]', text))

    counts = {'ko': korean, 'ja': japanese, 'zh': chinese, 'en': latin}
    dominant = max(counts, key=counts.get)
    if counts[dominant] == 0:
        return None
    return dominant


def get_font_for_language(lang):
    """Return font name based on detected language."""
    return {
        'ko': FONT_KOREAN,
        'ja': FONT_JAPANESE,
        'zh': FONT_CHINESE,
        'en': FONT_ENGLISH,
    }.get(lang, FONT_DEFAULT)



def clean_artifacts(prs):
    """Remove _x000B_ text artifacts left by translation model."""
    count = 0
    total = len(prs.slides)
    for pos, slide in enumerate(prs.slides):
        frames = list(_iter_all_frames(slide))
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            frames.append(slide.notes_slide.notes_text_frame)
        for frame in frames:
            for para in frame.paragraphs:
                for run in para.runs:
                    if '_x000B_' in run.text:
                        run.text = run.text.replace('_x000B_', '\n')
                        count += 1
        progress("Artifacts", pos + 1, total, f"Slide {pos+1}: {count} fixed")
    print()
    return count

def normalize_fonts(prs):
    """Replace fonts on all text runs based on content language,
    and patch XML directly to eliminate missing font references (e.g. NanumSquareOTF)."""

    count = 0

    # 1) Run-level: set font based on detected language
    total = len(prs.slides)
    for pos, slide in enumerate(prs.slides):
        frames = list(_iter_all_frames(slide))
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            frames.append(slide.notes_slide.notes_text_frame)

        for frame in frames:
            for para in frame.paragraphs:
                for run in para.runs:
                    if not run.text.strip():
                        continue
                    lang = detect_language(run.text)
                    if lang:
                        run.font.name = get_font_for_language(lang)
                        count += 1
        progress("Fonts", pos + 1, total, f"Slide {pos+1}: {count} runs")
    print()

    # 2) XML-level: replace all NanumSquare* typeface references across entire presentation
    xml_count = 0
    parts_to_patch = []

    # Collect all XML parts: slides, slide layouts, slide masters, theme, noteSlides
    for slide in prs.slides:
        parts_to_patch.append(slide._element)
        if slide.has_notes_slide:
            parts_to_patch.append(slide.notes_slide._element)
    for layout in prs.slide_layouts:
        parts_to_patch.append(layout._element)
    for master in prs.slide_masters:
        parts_to_patch.append(master._element)

    for part in parts_to_patch:
        for elem in part.iter():
            typeface = elem.get('typeface')
            if typeface and typeface in FONT_REPLACE:
                elem.set('typeface', FONT_REPLACE[typeface])
                xml_count += 1

    # 3) Patch docProps/app.xml (font list in package)
    try:
        from pptx.opc.constants import RELATIONSHIP_TYPE as RT
        app_part = prs.part.package.part_related_by(RT.CORE_PROPERTIES)
    except Exception:
        pass  # not critical

    return count, xml_count


def translate_batch(texts, model_id):
    """Translate a batch of texts from source to target language via Bedrock."""
    if not texts:
        return []

    src_name = LANG_NAMES.get(SOURCE_LANG, SOURCE_LANG)
    tgt_name = LANG_NAMES.get(TARGET_LANG, TARGET_LANG)

    prompt = (
        f"Translate the following {src_name} texts to natural, fluent {tgt_name}. "
        "Return ONLY a JSON array of translated strings, one per input, with no other text.\n\n"
        "Rules:\n"
        f"- Produce natural {tgt_name} sentence structure (do NOT translate word-by-word)\n"
        "- Localize date formats naturally (e.g. '2024년 10월' → 'October 2024', or 'Oct 2024' → '2024년 10월')\n"
        "- Keep technical terms, product names (Kiro, AWS, DynamoDB, SQS, etc.) as-is\n"
        "- For company names in non-Latin scripts, keep original + add romanization in parentheses\n"
        "- Preserve any markdown, bullet points, or line breaks in the original\n"
    )
    if GLOSSARY:
        prompt += "\n- MANDATORY glossary (use these exact translations):\n"
        for src, tgt in GLOSSARY.items():
            prompt += f"  \"{src}\" → \"{tgt}\"\n"
    prompt += f"\nInput ({len(texts)} items):\n" + json.dumps(texts, ensure_ascii=False)

    body = json.dumps({
        "anthropic_version": "bedrock-2023-05-31",
        "max_tokens": 16384,
        "temperature": 0.1,
        "messages": [{"role": "user", "content": prompt}]
    })

    for attempt in range(MAX_RETRIES):
        try:
            resp = client.invoke_model(modelId=model_id, body=body, contentType="application/json")
            resp_body = json.loads(resp['body'].read())
            # Track token usage
            global TOTAL_INPUT_TOKENS, TOTAL_OUTPUT_TOKENS
            usage = resp_body.get('usage', {})
            TOTAL_INPUT_TOKENS += usage.get('input_tokens', 0)
            TOTAL_OUTPUT_TOKENS += usage.get('output_tokens', 0)
            result = resp_body['content'][0]['text']
            parsed = _parse_json_array(result)
            if parsed and len(parsed) == len(texts):
                return parsed
        except Exception as e:
            if attempt < MAX_RETRIES - 1:
                time.sleep(2 ** attempt)
                continue
            print(f"\n  ⚠ API error after {MAX_RETRIES} retries: {e}", file=sys.stderr)

    return texts  # fallback: return originals


def _parse_json_array(text):
    """Robustly extract a JSON array from model response."""
    start = text.find('[')
    if start < 0:
        return None
    depth = 0
    for i in range(start, len(text)):
        if text[i] == '[':
            depth += 1
        elif text[i] == ']':
            depth -= 1
        if depth == 0:
            try:
                return json.loads(text[start:i + 1])
            except json.JSONDecodeError:
                return None
    return None


def collect_translatable_paragraphs(frame):
    """Collect paragraphs containing source language text from a text frame."""
    entries = []
    for pi, para in enumerate(frame.paragraphs):
        if para.text.strip() and has_source_lang(para.text):
            entries.append((pi, para.text))
    return entries


def apply_translation(frame, para_idx, translated_text):
    """Apply translated text to paragraph: first run gets full text, rest cleared."""
    para = frame.paragraphs[para_idx]
    if not para.runs:
        return
    para.runs[0].text = translated_text
    for run in para.runs[1:]:
        run.text = ""


def _iter_all_frames(slide):
    """Yield all text frames from a slide, including groups and tables."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    def _from_shape(shape):
        if shape.has_text_frame:
            yield shape.text_frame
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    yield cell.text_frame
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                yield from _from_shape(child)

    for shape in slide.shapes:
        yield from _from_shape(shape)


def progress(label, current, total, detail=""):
    """Display real-time progress bar."""
    pct = current / total
    bar_len = 30
    filled = int(bar_len * pct)
    bar = '█' * filled + '░' * (bar_len - filled)
    sys.stdout.write(f"\r  {label} [{bar}] {current}/{total}  {detail}    ")
    sys.stdout.flush()


def translate_pass(prs, model_id, label="Translating", slide_indices=None, include_notes=True):
    """Single translation pass over all (or selected) slides with parallel API calls."""
    from concurrent.futures import ThreadPoolExecutor, as_completed
    import threading

    total_slides = len(prs.slides)
    total_translated = 0
    indices = list(slide_indices if slide_indices is not None else range(total_slides))
    process_total = len(indices)
    progress_lock = threading.Lock()
    completed = [0]  # mutable counter for threads

    def process_slide(si):
        """Collect and translate paragraphs for one slide. Returns (si, [(frame, pi, new_text)])."""
        slide = prs.slides[si]
        entries = []

        for frame in _iter_all_frames(slide):
            for pi, text in collect_translatable_paragraphs(frame):
                entries.append((frame, pi, text))

        if include_notes and slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            ntf = slide.notes_slide.notes_text_frame
            for pi, text in collect_translatable_paragraphs(ntf):
                entries.append((ntf, pi, text))

        if not entries:
            return si, []

        texts = [e[2] for e in entries]
        translated = []
        for i in range(0, len(texts), BATCH_SIZE):
            translated.extend(translate_batch(texts[i:i + BATCH_SIZE], model_id))

        return si, list(zip(entries, translated))

    # Process slides in parallel (max 4 concurrent API calls)
    results = {}
    with ThreadPoolExecutor(max_workers=4) as executor:
        futures = {executor.submit(process_slide, si): si for si in indices}
        for future in as_completed(futures):
            si = futures[future]
            si_result, translations = future.result()
            results[si] = translations
            with progress_lock:
                completed[0] += 1
                count = len(translations)
                progress(label, completed[0], process_total,
                         f"Slide {si+1}: {count} paragraphs" if count else f"Slide {si+1}: skip")

    # Apply translations in slide order (must be sequential for python-pptx)
    for si in indices:
        for (frame, pi, _), new_text in results.get(si, []):
            apply_translation(frame, pi, new_text)
            total_translated += 1

    print()
    return total_translated


def review_pass(prs, slide_indices=None, include_notes=True):
    """Scan all (or selected) slides and return list of (slide_num, location, text) with remaining source language."""
    remaining = []
    indices = slide_indices if slide_indices is not None else range(len(prs.slides))
    total = len(indices) if slide_indices is not None else len(prs.slides)
    for pos, i in enumerate(indices):
        slide = prs.slides[i]
        for frame in _iter_all_frames(slide):
            for para in frame.paragraphs:
                if para.text.strip() and has_source_lang(para.text):
                    remaining.append((i + 1, 'text', para.text[:120]))
        if include_notes and slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            for para in slide.notes_slide.notes_text_frame.paragraphs:
                if para.text.strip() and has_source_lang(para.text):
                    remaining.append((i + 1, 'note', para.text[:120]))
        progress("Review", pos + 1, total, f"Slide {i+1}: {len(remaining)} found")
    print()
    return remaining



def summarize_presentation(prs, model_id, slide_indices=None):
    """Extract all text from slides and generate a structured summary via Bedrock."""
    total = len(prs.slides)
    indices = slide_indices if slide_indices is not None else range(total)
    slide_contents = []
    idx_list = list(indices)

    for pos, si in enumerate(idx_list):
        slide = prs.slides[si]
        texts = []
        for frame in _iter_all_frames(slide):
            t = frame.text.strip()
            if t:
                texts.append(t)
        notes = ''
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        if texts or notes:
            entry = f"--- Slide {si+1} ---\n"
            if texts:
                entry += "Content: " + " | ".join(texts) + "\n"
            if notes:
                entry += "Notes: " + notes + "\n"
            slide_contents.append(entry)
        progress("Extract", pos + 1, len(idx_list), f"Slide {si+1}")
    print()

    all_text = "\n".join(slide_contents)
    # Truncate if too long for a single API call
    if len(all_text) > 80000:
        all_text = all_text[:80000] + "\n...(truncated)"

    prompt = (
        "Analyze this presentation and provide a structured summary:\n\n"
        "1. **Title & Speaker** — presentation title and presenter\n"
        "2. **Executive Summary** — 3-5 sentence overview\n"
        "3. **Key Topics** — bullet list of main topics covered\n"
        "4. **Slide-by-Slide Outline** — one-line summary per slide (group related slides)\n"
        "5. **Key Takeaways** — the most important points for the audience\n\n"
        "Presentation content:\n\n" + all_text
    )

    body = json.dumps({
        "anthropic_version": "bedrock-2023-05-31",
        "max_tokens": 8192,
        "temperature": 0.2,
        "messages": [{"role": "user", "content": prompt}]
    })

    sys.stdout.write("\r  🤖 Generating summary with Bedrock...    ")
    sys.stdout.flush()
    resp = client.invoke_model(modelId=model_id, body=body, contentType="application/json")
    sys.stdout.write("\r                                           \r")
    return json.loads(resp['body'].read())['content'][0]['text']

def parse_slide_range(slide_str, total):
    """Parse slide range string like '1-10' or '1,3,5-8' into a set of 0-based indices."""
    indices = set()
    for part in slide_str.split(','):
        part = part.strip()
        if '-' in part:
            a, b = part.split('-', 1)
            indices.update(range(int(a) - 1, int(b)))
        else:
            indices.add(int(part) - 1)
    return sorted(i for i in indices if 0 <= i < total)


def dry_run(prs, slide_indices=None, include_notes=True):
    """Preview translation scope without making API calls."""
    indices = slide_indices if slide_indices is not None else range(len(prs.slides))
    total_paras = 0
    total_chars = 0
    slides_with_content = 0

    for si in indices:
        slide = prs.slides[si]
        slide_paras = 0
        for frame in _iter_all_frames(slide):
            for pi, text in collect_translatable_paragraphs(frame):
                slide_paras += 1
                total_chars += len(text)
        if include_notes and slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            for pi, text in collect_translatable_paragraphs(slide.notes_slide.notes_text_frame):
                slide_paras += 1
                total_chars += len(text)
        if slide_paras:
            slides_with_content += 1
        total_paras += slide_paras

    api_calls = (total_paras + BATCH_SIZE - 1) // BATCH_SIZE if total_paras else 0

    print(f"{'='*60}")
    print(f"📊 Dry Run — Translation Preview")
    print(f"{'='*60}")
    print(f"   📄 Slides to process:  {len(list(indices))}")
    print(f"   📝 Slides with content: {slides_with_content}")
    print(f"   📋 Paragraphs to translate: {total_paras}")
    print(f"   🔤 Total characters: {total_chars:,}")
    print(f"   🌐 API calls (batch {BATCH_SIZE}): {api_calls}")
    print(f"   📌 Notes included: {'Yes' if include_notes else 'No'}")
    if GLOSSARY:
        print(f"   📖 Glossary terms: {len(GLOSSARY)}")
    # Rough cost estimate: ~4 chars per token for mixed CJK/Latin
    est_input_tokens = total_chars // 2  # prompt overhead ~2x content
    est_output_tokens = int(est_input_tokens * 1.5)  # translations are longer
    est_cost = (est_input_tokens * 15 + est_output_tokens * 75) / 1_000_000
    print(f"   💰 Est. cost: ~${est_cost:.3f} (Opus pricing, rough estimate)")


def save_review_report(prs_orig, prs_trans, output_path, slide_indices=None):
    """Generate a markdown comparison report of original vs translated notes."""
    indices = slide_indices if slide_indices is not None else range(len(prs_orig.slides))
    lines = ["# Translation Review Report\n"]
    count = 0

    for i in indices:
        o_notes = ''
        t_notes = ''
        if prs_orig.slides[i].has_notes_slide and prs_orig.slides[i].notes_slide.notes_text_frame:
            o_notes = prs_orig.slides[i].notes_slide.notes_text_frame.text.strip()
        if prs_trans.slides[i].has_notes_slide and prs_trans.slides[i].notes_slide.notes_text_frame:
            t_notes = prs_trans.slides[i].notes_slide.notes_text_frame.text.strip()
        if not o_notes and not t_notes:
            continue
        count += 1
        ratio = f"{len(t_notes)/len(o_notes):.0%}" if o_notes else "N/A"
        lines.append(f"## Slide {i+1} (length ratio: {ratio})\n")
        lines.append(f"### Original\n{o_notes}\n")
        lines.append(f"### Translated\n{t_notes}\n")
        lines.append("---\n")

    with open(output_path, 'w', encoding='utf-8') as f:
        f.write('\n'.join(lines))
    return count


def main():
    parser = argparse.ArgumentParser(description='Translate PPTX between languages using AWS Bedrock')
    parser.add_argument('input', help='Input PPTX file path')
    parser.add_argument('-o', '--output', help='Output file path (default: <input>_<target>.pptx)')
    parser.add_argument('-s', '--source', default='ko', help='Source language code (default: ko)')
    parser.add_argument('-t', '--target', default='en', help='Target language code (default: en)')
    parser.add_argument('--mode', default='all', choices=['translate', 'review', 'fonts', 'summary', 'all'],
                        help='Execution mode (default: all)')
    parser.add_argument('--slides', help='Slide range to process (e.g. "1-10" or "1,3,5-8")')
    parser.add_argument('--model', default=MODEL_ID, help=f'Bedrock model ID (default: {MODEL_ID})')
    parser.add_argument('--max-review-passes', type=int, default=2, help='Max review/fix passes (default: 2)')
    parser.add_argument('--dry-run', action='store_true', help='Preview translation scope without API calls')
    parser.add_argument('--no-notes', action='store_true', help='Skip slide notes translation')
    parser.add_argument('--glossary', help='JSON glossary file for fixed term translations')
    parser.add_argument('--resume', action='store_true', help='Skip already-translated slides (no source language text)')
    parser.add_argument('--report', help='Output markdown review report path (used with --mode review)')
    parser.add_argument('--original', help='Original PPTX for comparison report (used with --report)')
    args = parser.parse_args()

    global SOURCE_LANG, TARGET_LANG, GLOSSARY
    SOURCE_LANG = args.source
    TARGET_LANG = args.target

    # Load glossary
    if args.glossary:
        with open(args.glossary, 'r', encoding='utf-8') as f:
            GLOSSARY = json.load(f)
        print(f"📖 Glossary loaded: {len(GLOSSARY)} terms")

    src_name = LANG_NAMES.get(SOURCE_LANG, SOURCE_LANG)
    tgt_name = LANG_NAMES.get(TARGET_LANG, TARGET_LANG)
    tgt_suffix = LANG_NAMES.get(TARGET_LANG, TARGET_LANG).lower()
    output = args.output or args.input.replace('.pptx', f'_{tgt_suffix}.pptx')

    # For review/fonts mode on already-translated file, default output = input
    if args.mode in ('review', 'fonts') and not args.output:
        output = args.input

    print(f"📂 Input:  {args.input}")
    print(f"📂 Output: {output}")
    print(f"🌐 {src_name} → {tgt_name}")
    print(f"🔧 Mode:   {args.mode}")
    if args.slides:
        print(f"📄 Slides: {args.slides}")
    print(f"🤖 Model:  {args.model}")
    print()

    prs = Presentation(args.input)
    total_slides = len(prs.slides)
    slide_indices = parse_slide_range(args.slides, total_slides) if args.slides else None
    include_notes = not args.no_notes
    start_time = time.time()

    # ── Dry run ──
    if args.dry_run:
        dry_run(prs, slide_indices, include_notes)
        return

    # ── Resume: filter out already-translated slides ──
    if args.resume and args.mode in ('translate', 'all'):
        all_indices = slide_indices if slide_indices is not None else list(range(total_slides))
        remaining_indices = []
        for si in all_indices:
            slide = prs.slides[si]
            has_src = False
            for frame in _iter_all_frames(slide):
                for para in frame.paragraphs:
                    if para.text.strip() and has_source_lang(para.text):
                        has_src = True
                        break
                if has_src:
                    break
            if not has_src and include_notes and slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                for para in slide.notes_slide.notes_text_frame.paragraphs:
                    if para.text.strip() and has_source_lang(para.text):
                        has_src = True
                        break
            if has_src:
                remaining_indices.append(si)
        skipped = len(all_indices) - len(remaining_indices)
        if skipped:
            print(f"⏩ Resume: skipping {skipped} already-translated slides")
        slide_indices = remaining_indices if remaining_indices else None
        if not remaining_indices:
            print("   ✅ All slides already translated!")
            return

    # ── Summary ──
    if args.mode == 'summary':
        print("📝 Generating presentation summary...")
        summary = summarize_presentation(prs, args.model, slide_indices)
        print(summary)
        print(f"\n✅ Summary complete")
        return

    # ── Translate ──
    if args.mode in ('translate', 'all'):
        print("🔄 Translating...")
        count = translate_pass(prs, args.model, "Translate", slide_indices, include_notes)
        print(f"   ✅ {count} paragraphs translated")

        # Review & fix passes
        for pass_num in range(2, 2 + args.max_review_passes):
            print(f"\n🔍 Review: Scanning for remaining {src_name} text...")
            remaining = review_pass(prs, slide_indices, include_notes)

            if not remaining:
                print("   ✅ No remaining source language text found!")
                break

            real_issues = [r for r in remaining if not re.search(r'[A-Za-z].*\([^)]*[가-힣ぁ-ヶ一-鿿]|[가-힣ぁ-ヶ一-鿿].*\(.*[A-Za-z]', r[2])]

            if not real_issues:
                print(f"   ✅ {len(remaining)} item(s) with intentional {src_name} (names with romanization) — OK")
                break

            slides_affected = sorted(set(r[0] for r in real_issues))
            print(f"   ⚠ {len(real_issues)} untranslated paragraph(s) in slides: {slides_affected}")

            print(f"\n🔄 Fix pass {pass_num}...")
            count = translate_pass(prs, args.model, f"Fix {pass_num}", slide_indices, include_notes)
            print(f"   ✅ {count} paragraphs fixed")

    # ── Review only ──
    if args.mode == 'review':
        print("🔍 Reviewing translation quality...")
        remaining = review_pass(prs, slide_indices, include_notes)
        if not remaining:
            print("   ✅ No source language text found!")
        else:
            real_issues = [r for r in remaining if not re.search(r'[A-Za-z].*\([^)]*[가-힣ぁ-ヶ一-鿿]|[가-힣ぁ-ヶ一-鿿].*\(.*[A-Za-z]', r[2])]
            intentional = len(remaining) - len(real_issues)
            if intentional:
                print(f"   ✅ {intentional} intentional (names with romanization)")
            if real_issues:
                print(f"   ⚠  {len(real_issues)} untranslated paragraph(s):")
                for slide_num, loc, text in real_issues:
                    print(f"      Slide {slide_num} [{loc}]: {text[:80]}")
            else:
                print("   ✅ All translations look good!")

        # Generate comparison report if requested
        if args.report and args.original:
            prs_orig = Presentation(args.original)
            report_count = save_review_report(prs_orig, prs, args.report, slide_indices)
            print(f"\n📄 Review report saved: {args.report} ({report_count} slides compared)")
        elif args.report:
            print(f"\n⚠  --original required for comparison report (provide original PPTX path)")

    # ── Fonts only ──
    if args.mode in ('fonts', 'all'):
        print("🧹 Cleaning translation artifacts...")
        artifact_count = clean_artifacts(prs)
        print(f"   ✅ {artifact_count} artifacts fixed" if artifact_count else "   ✅ No artifacts found")

        print("🔤 Normalizing fonts...")
        run_count, xml_count = normalize_fonts(prs)
        print(f"   ✅ {run_count} text runs updated (language-based)")
        print(f"   ✅ {xml_count} XML font references patched")
        print(f"      Korean → {FONT_KOREAN} | English → {FONT_ENGLISH}")

    # ── Final summary (all mode) ──
    if args.mode == 'all':
        remaining = review_pass(prs, None, include_notes)
        elapsed = time.time() - start_time
        print(f"\n{'='*60}")
        print(f"📊 Translation Summary")
        print(f"{'='*60}")
        print(f"   ⏱  Time: {elapsed:.0f}s")
        if TOTAL_INPUT_TOKENS or TOTAL_OUTPUT_TOKENS:
            # Opus pricing: $15/M input, $75/M output
            cost = (TOTAL_INPUT_TOKENS * 15 + TOTAL_OUTPUT_TOKENS * 75) / 1_000_000
            print(f"   🔢 Tokens: {TOTAL_INPUT_TOKENS:,} input + {TOTAL_OUTPUT_TOKENS:,} output")
            print(f"   💰 Est. cost: ${cost:.3f} (Opus pricing)")
        if remaining:
            intentional = [r for r in remaining if re.search(r'[A-Za-z].*\([^)]*[가-힣ぁ-ヶ一-鿿]|[가-힣ぁ-ヶ一-鿿].*\(.*[A-Za-z]', r[2])]
            issues = [r for r in remaining if r not in intentional]
            if intentional:
                print(f"   ✅ {len(intentional)} intentional {src_name} (names with romanization)")
            if issues:
                print(f"   ⚠  {len(issues)} paragraph(s) still contain {src_name}:")
                for slide_num, loc, text in issues:
                    print(f"      Slide {slide_num} [{loc}]: {text[:80]}")
        else:
            print("   ✅ All source language text successfully translated!")

    # ── Save ──
    if args.mode != 'review':
        print(f"\n💾 Saving to {output}...")
        prs.save(output)

        import zipfile, shutil
        tmp_path = output + '.tmp'
        with zipfile.ZipFile(output, 'r') as zin, zipfile.ZipFile(tmp_path, 'w') as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename == 'docProps/app.xml':
                    text = data.decode('utf-8')
                    for old_font in FONT_REPLACE:
                        text = text.replace(old_font, FONT_REPLACE[old_font])
                    data = text.encode('utf-8')
                zout.writestr(item, data)
        shutil.move(tmp_path, output)
        print(f"✅ Done!")
    else:
        print(f"\n✅ Review complete (no changes saved)")


if __name__ == '__main__':
    main()
